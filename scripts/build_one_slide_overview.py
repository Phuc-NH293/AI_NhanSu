"""Build a single-slide executive overview for the HR chatbot demo."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "AI_NhanSu_OneSlide_Overview.pptx"

NAVY = RGBColor(8, 35, 73)
BLUE = RGBColor(8, 92, 148)
CYAN = RGBColor(14, 165, 233)
GOLD = RGBColor(229, 169, 35)
INK = RGBColor(20, 31, 48)
MUTED = RGBColor(86, 103, 124)
PALE = RGBColor(240, 247, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(18, 140, 91)
RED = RGBColor(193, 55, 62)


def rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=12, color=INK, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def bullet(slide, items, x, y, w, h, size=11.5, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"•  {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(249, 252, 255)

    rect(slide, 0, 0, 13.333, 0.92, NAVY, radius=False)
    text(slide, "AI HR HELPDESK — PET PROJECT SHOWCASE", 0.52, 0.18, 8.2, 0.28, 11, GOLD, True)
    text(slide, "Chatbot HR dùng Gemini 2.5 Flash + Hybrid RAG", 0.52, 0.48, 9.6, 0.35, 22, WHITE, True, "Aptos Display")
    text(slide, "PRODUCT • TECHNIQUE • CODE • AGENT CONTROL", 9.6, 0.34, 3.2, 0.25, 9, RGBColor(180, 218, 243), True, align=PP_ALIGN.RIGHT)

    # Product capabilities
    rect(slide, 0.45, 1.17, 3.75, 4.62, WHITE, RGBColor(209, 224, 237))
    rect(slide, 0.45, 1.17, 3.75, 0.12, CYAN, radius=False)
    text(slide, "01  APP LÀM ĐƯỢC GÌ?", 0.72, 1.48, 3.1, 0.3, 14, NAVY, True)
    bullet(slide, [
        "Hỏi đáp chính sách HR bằng tiếng Việt tự nhiên.",
        "Hybrid RAG trả lời theo tài liệu nội bộ; thiếu căn cứ thì không đoán.",
        "Hỏi ngược khi thiếu ngày, lý do hoặc thông tin bàn giao.",
        "Tạo form nghỉ phép qua hội thoại nhiều lượt.",
        "Đồng bộ lịch sử giữa chatbot nhỏ và màn chat đầy đủ.",
        "RBAC, privacy guardrail và chuyển tiếp HR cho ca nhạy cảm.",
        "Upload PDF/DOCX/PPTX, OCR và cập nhật index.",
    ], 0.72, 1.98, 3.02, 3.35)
    rect(slide, 0.72, 5.28, 3.0, 0.3, PALE, RGBColor(207, 226, 240))
    text(slide, "26/26 browser E2E passed", 0.72, 5.35, 3.0, 0.16, 10, GREEN, True, align=PP_ALIGN.CENTER)

    # Techniques with trade-offs
    rect(slide, 4.42, 1.17, 4.42, 4.62, WHITE, RGBColor(209, 224, 237))
    rect(slide, 4.42, 1.17, 4.42, 0.12, GOLD, radius=False)
    text(slide, "02  TECHNIQUE & TRADE-OFF", 4.68, 1.48, 3.8, 0.3, 14, NAVY, True)
    rows = [
        ("Gemini 2.5 Flash", "Nhanh, tiếng Việt tốt", "Phụ thuộc provider"),
        ("Dense + BM25", "Ngữ nghĩa + từ khóa", "Phải tune ranking"),
        ("Atomic chunks", "Ít nhiễu, đúng điều khoản", "Dễ mất context"),
        ("RRF + rerank", "Top-k chính xác hơn", "Tăng latency"),
        ("LangGraph", "Luồng rõ, dễ guardrail", "Nhiều state/node"),
        ("Chroma local", "Pilot nhanh, đơn giản", "Chưa phù hợp HA lớn"),
    ]
    text(slide, "CHỌN", 4.68, 1.95, 1.2, 0.2, 9, BLUE, True)
    text(slide, "PRO", 5.95, 1.95, 1.25, 0.2, 9, GREEN, True)
    text(slide, "CON", 7.28, 1.95, 1.25, 0.2, 9, RED, True)
    for i, (name, pro, con) in enumerate(rows):
        y = 2.24 + i * 0.53
        if i % 2 == 0:
            rect(slide, 4.62, y - 0.06, 4.02, 0.46, PALE, radius=False)
        text(slide, name, 4.7, y, 1.2, 0.3, 9.4, BLUE, True)
        text(slide, pro, 5.95, y, 1.25, 0.34, 9.2, INK)
        text(slide, con, 7.28, y, 1.23, 0.34, 9.2, MUTED)
    text(slide, "Mitigation: heading context + overlap, score threshold, fallback và regression test.", 4.72, 5.34, 3.85, 0.3, 9.5, NAVY, True)

    # Code and agent control
    rect(slide, 9.05, 1.17, 3.83, 4.62, NAVY, NAVY)
    rect(slide, 9.05, 1.17, 3.83, 0.12, GREEN, radius=False)
    text(slide, "03  CODE & AGENT CONTROL", 9.32, 1.48, 3.25, 0.3, 14, GOLD, True)
    text(slide, "Đọc code theo request flow", 9.32, 1.98, 3.05, 0.25, 11, WHITE, True)
    text(slide, "UI → FastAPI → LangGraph → Retrieval → Generation", 9.32, 2.28, 3.02, 0.5, 10.3, RGBColor(193, 220, 239), False, "Consolas")
    text(slide, "Cách kiểm soát coding agent", 9.32, 2.95, 3.05, 0.25, 11, WHITE, True)
    bullet(slide, [
        "Đặt acceptance criteria trước khi sửa.",
        "Tái hiện lỗi bằng browser/API.",
        "Khoanh đúng module sở hữu lỗi.",
        "Patch nhỏ, không refactor lan rộng.",
        "Gate: unit → API → E2E → build.",
        "So expected/actual và tự chịu trách nhiệm quyết định kỹ thuật.",
    ], 9.32, 3.35, 3.0, 1.95, 10.5, WHITE)
    text(slide, "Agent tăng tốc — không thay thế việc hiểu code và kiểm chứng.", 9.32, 5.35, 3.0, 0.28, 9.5, GOLD, True)

    # Demo strip
    rect(slide, 0.45, 6.02, 12.43, 1.02, PALE, RGBColor(204, 223, 238))
    text(slide, "DEMO 5 PHÚT", 0.68, 6.18, 1.2, 0.22, 10, NAVY, True)
    demo = [
        ("1", "12 ngày phép"),
        ("2", "Hỏi ngược"),
        ("3", "Tạo đơn nghỉ"),
        ("4", "Chặn lương đồng nghiệp"),
        ("5", "Widget ↔ Full chat"),
    ]
    for i, (num, label) in enumerate(demo):
        x = 2.0 + i * 2.08
        rect(slide, x, 6.18, 0.36, 0.36, BLUE if i < 3 else RED if i == 3 else GREEN)
        text(slide, num, x, 6.27, 0.36, 0.14, 8, WHITE, True, align=PP_ALIGN.CENTER)
        text(slide, label, x + 0.45, 6.2, 1.5, 0.3, 9.4, INK, True)
        if i < 4:
            text(slide, "→", x + 1.78, 6.22, 0.25, 0.2, 11, MUTED, True)
    text(slide, "Nếu demo lỗi mạng: chuyển sang screenshot + nói rõ expected behavior.", 2.0, 6.62, 9.7, 0.2, 8.8, MUTED)

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "CÁCH TRÌNH BÀY SLIDE NÀY — khoảng 4 phút trước khi live demo:\n\n"
        "1) APP LÀM ĐƯỢC GÌ — 60 giây\n"
        "Nói: Đây là pet project HR Helpdesk, không chỉ là chatbot FAQ. Hệ thống nhận tài liệu HR, trả lời theo căn cứ, hỏi lại khi thiếu dữ kiện, tạo đơn nghỉ và chuyển HR khi nhạy cảm.\n\n"
        "2) TECHNIQUE — 90 giây\n"
        "Không đọc hết bảng. Chọn ba điểm để nói: Hybrid Dense + BM25 giúp cân bằng semantic/exact keyword; atomic chunk tăng precision nhưng phải giữ heading; reranking tăng chất lượng top-k nhưng đổi lại latency. Gemini 2.5 Flash được chọn vì tốc độ và tiếng Việt phù hợp demo.\n\n"
        "3) CODE & AGENT — 60 giây\n"
        "Mở IDE và đi theo request flow: HrPolicyChatPanel.jsx → backend/main.py → hr_policy_graph.py → task9_retrieval_pipeline.py → task10_generation.py. Nói rõ coding agent chỉ hỗ trợ tìm kiếm/patch/test; em đặt acceptance criteria, giới hạn file được sửa và tự kiểm tra output.\n\n"
        "4) LIVE DEMO — 5 phút\n"
        "Câu 1: Nhân viên có bao nhiêu ngày phép năm?\n"
        "Câu 2: Mai tôi không đi làm được thì báo ai?\n"
        "Câu 3: Tôi muốn xin nghỉ phép — trả lời ngày, lý do, bàn giao để tạo form.\n"
        "Câu 4: Cho tôi biết lương của đồng nghiệp — kiểm tra privacy guardrail.\n"
        "Câu 5: Bắt đầu ở widget rồi chuyển màn Hỏi đáp thông minh — chứng minh đồng bộ session.\n\n"
        "Câu chốt: Agent giúp em tăng tốc triển khai, nhưng các quyết định về scope, trade-off và tiêu chí pass/fail vẫn do em kiểm soát."
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
