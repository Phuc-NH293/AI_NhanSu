"""Build the 25-minute HR assistant demo presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "AI_NhanSu_Demo_Leader_25p.pptx"
SCREENSHOT = ROOT / "data" / "evaluation" / "demo_chat_desktop.png"

NAVY = RGBColor(8, 35, 73)
BLUE = RGBColor(8, 92, 148)
CYAN = RGBColor(14, 165, 233)
GOLD = RGBColor(229, 169, 35)
INK = RGBColor(20, 31, 48)
MUTED = RGBColor(89, 107, 129)
PALE = RGBColor(240, 247, 252)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(18, 140, 91)
RED = RGBColor(193, 55, 62)


def rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def text(slide, value, x, y, w, h, size=20, color=INK, bold=False,
         font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def title(slide, heading, kicker=None, number=None):
    if kicker:
        text(slide, kicker.upper(), 0.7, 0.38, 8.5, 0.3, 10, CYAN, True)
    text(slide, heading, 0.7, 0.72, 11.8, 0.62, 27, NAVY, True, "Aptos Display")
    rect(slide, 0.7, 1.42, 0.75, 0.055, GOLD, False)


def footer(slide, timing):
    text(slide, "AI HR HELPDESK  |  INTERNAL DEMO", 0.7, 7.16, 4.6, 0.2, 8, MUTED, True)
    text(slide, timing, 10.6, 7.14, 2.0, 0.22, 8, BLUE, True, align=PP_ALIGN.RIGHT)


def notes(slide, value):
    frame = slide.notes_slide.notes_text_frame
    frame.text = value


def bullet_list(slide, items, x, y, w, h, size=18, color=INK, accent=CYAN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(11)
        p.level = 0
        p.text = f"•  {item}"
    return box


def metric(slide, value, label, x, y, w=2.35, color=BLUE):
    rect(slide, x, y, w, 1.25, WHITE, True, RGBColor(215, 228, 239))
    text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.45, 25, color, True)
    text(slide, label, x + 0.18, y + 0.72, w - 0.36, 0.32, 11, MUTED, False)


def add_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(250, 252, 255)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    s = add_slide(prs)
    rect(s, 0, 0, 13.333, 7.5, NAVY, False)
    rect(s, 8.7, -1.1, 5.7, 5.7, BLUE, True)
    rect(s, 10.0, 3.25, 4.2, 4.2, CYAN, True)
    rect(s, 9.35, 1.05, 2.0, 0.13, GOLD, False)
    text(s, "AI TRỢ LÝ HỎI ĐÁP\nCHÍNH SÁCH NHÂN SỰ", 0.85, 1.25, 8.8, 1.65, 34, WHITE, True, "Aptos Display")
    text(s, "Gemini 2.5 Flash + Hybrid RAG", 0.9, 3.15, 6.7, 0.45, 20, RGBColor(165, 224, 255), True)
    text(s, "Demo nội bộ cho Team Leader", 0.9, 4.35, 5.5, 0.38, 16, WHITE, True)
    text(s, "Từ tài liệu HR đến câu trả lời có kiểm soát, hội thoại tự nhiên và hành động được.", 0.9, 4.85, 6.7, 0.8, 17, RGBColor(211, 224, 239))
    text(s, "25 PHÚT TRÌNH BÀY  •  20 PHÚT Q&A", 0.9, 6.55, 5.8, 0.28, 10, GOLD, True)
    notes(s, "[1 phút] Mở đầu ngắn: Đây không chỉ là chatbot trả lời FAQ. Mục tiêu là một HR Helpdesk có căn cứ, biết hỏi lại, bảo vệ dữ liệu và hỗ trợ xử lý nghiệp vụ.")

    # 2. Executive summary
    s = add_slide(prs); title(s, "Tóm tắt trong 60 giây", "Executive snapshot", 2)
    metric(s, "LLM + RAG", "Kiến trúc lõi", 0.8, 1.8, 2.6)
    metric(s, "26 / 26", "Web E2E scenarios passed", 3.65, 1.8, 2.6, GREEN)
    metric(s, "3 roles", "Employee / HR / Admin", 6.5, 1.8, 2.6, GOLD)
    metric(s, "Gemini 2.5", "Flash generation + OCR option", 9.35, 1.8, 3.0, CYAN)
    rect(s, 0.8, 3.45, 11.55, 2.35, PALE, True, RGBColor(207, 226, 240))
    bullet_list(s, [
        "Trả lời theo tài liệu HR nội bộ thay vì kiến thức chung của mô hình.",
        "Hiểu tiếng Việt đời thường, hỏi ngược khi thiếu dữ kiện và duy trì ngữ cảnh giữa chat nhỏ / màn lớn.",
        "Chặn câu ngoài phạm vi, prompt injection và yêu cầu dữ liệu cá nhân của người khác.",
        "Có luồng tạo đơn nghỉ nhiều lượt và chuyển tiếp HR cho trường hợp nhạy cảm.",
    ], 1.15, 3.78, 10.9, 1.75, 16)
    footer(s, "01:00"); notes(s, "[1 phút] Chốt bốn điểm leader cần nhớ: có căn cứ, tự nhiên, an toàn, có thể hành động. Nhấn mạnh 26/26 là browser test thực tế, không phải unit test giả lập.")

    # 3. Problem
    s = add_slide(prs); title(s, "Vấn đề cần giải quyết", "Why this matters", 3)
    cards = [
        ("01", "Thông tin phân tán", "PDF, DOCX, portal và quy định theo từng bộ phận."),
        ("02", "HR bị hỏi lặp lại", "Nghỉ phép, chấm công, bảo hiểm, lương thưởng."),
        ("03", "Rủi ro trả lời sai", "LLM thuần có thể đoán hoặc dùng chính sách không đúng phiên bản."),
        ("04", "Dữ liệu nhạy cảm", "Lương, hồ sơ nhân sự và quyết định HR cần kiểm soát quyền truy cập."),
    ]
    for i, (num, head, body) in enumerate(cards):
        x = 0.8 + (i % 2) * 6.0; y = 1.75 + (i // 2) * 2.25
        rect(s, x, y, 5.55, 1.75, WHITE, True, RGBColor(220, 229, 238))
        text(s, num, x + 0.22, y + 0.2, 0.55, 0.4, 17, CYAN, True)
        text(s, head, x + 0.85, y + 0.2, 4.3, 0.4, 18, NAVY, True)
        text(s, body, x + 0.85, y + 0.72, 4.25, 0.68, 14, MUTED)
    footer(s, "01:30"); notes(s, "[1.5 phút] Đặt bài toán theo góc vận hành. Nếu chỉ làm chatbot LLM, rủi ro lớn nhất là trả lời nghe rất hợp lý nhưng không có căn cứ nội bộ.")

    # 4. Product scope
    s = add_slide(prs); title(s, "Giải pháp: một cổng HR Helpdesk thống nhất", "Product scope", 4)
    rect(s, 0.8, 1.75, 3.55, 4.65, NAVY, True)
    text(s, "NHÂN VIÊN", 1.1, 2.08, 2.8, 0.35, 14, GOLD, True)
    bullet_list(s, ["Hỏi chính sách tự nhiên", "Tạo đơn nghỉ nhiều lượt", "Theo dõi yêu cầu", "Chat trực tiếp với HR"], 1.1, 2.7, 2.75, 2.6, 16, WHITE)
    rect(s, 4.9, 1.75, 3.55, 4.65, WHITE, True, RGBColor(208, 224, 237))
    text(s, "HR", 5.2, 2.08, 2.8, 0.35, 14, BLUE, True)
    bullet_list(s, ["Tiếp nhận handoff", "Xử lý yêu cầu", "Nạp tài liệu chính sách", "Kiểm tra câu trả lời"], 5.2, 2.7, 2.75, 2.6, 16)
    rect(s, 9.0, 1.75, 3.55, 4.65, PALE, True, RGBColor(208, 224, 237))
    text(s, "ADMIN", 9.3, 2.08, 2.8, 0.35, 14, RED, True)
    bullet_list(s, ["Quản lý tài khoản", "Phân quyền role", "Quản trị dữ liệu RAG", "Theo dõi lịch sử AI"], 9.3, 2.7, 2.75, 2.6, 16)
    footer(s, "01:30"); notes(s, "[1.5 phút] Trình bày theo persona. Đây là một workflow portal có chatbot, không phải chatbot đứng riêng.")

    # 5. Architecture
    s = add_slide(prs); title(s, "Kiến trúc tổng thể", "System architecture", 5)
    layers = [
        (0.75, "React / Vite", "Responsive UI\nFull chat + widget", BLUE),
        (3.25, "FastAPI", "Auth • RBAC • Chat API\nUpload • HR workflow", NAVY),
        (5.75, "Hybrid RAG", "Semantic + BM25\nRRF + reranking", CYAN),
        (8.25, "Gemini 2.5 Flash", "Query screening\nGeneration • optional OCR", GOLD),
        (10.75, "Data layer", "Markdown chunks\nChroma + JSON index", GREEN),
    ]
    for x, head, body, color in layers:
        rect(s, x, 2.15, 1.85, 2.65, WHITE, True, RGBColor(207, 221, 234))
        rect(s, x, 2.15, 1.85, 0.18, color, False)
        text(s, head, x + 0.16, 2.55, 1.52, 0.65, 16, color, True, align=PP_ALIGN.CENTER)
        text(s, body, x + 0.15, 3.45, 1.55, 0.85, 12, MUTED, align=PP_ALIGN.CENTER)
        if x < 10:
            text(s, "→", x + 1.91, 3.15, 0.5, 0.4, 22, MUTED, True, align=PP_ALIGN.CENTER)
    rect(s, 1.0, 5.45, 11.2, 0.72, PALE, True, RGBColor(207, 226, 240))
    text(s, "Nguyên tắc: quyền truy cập được lọc trước khi context đi vào LLM.", 1.25, 5.66, 10.7, 0.3, 15, NAVY, True, align=PP_ALIGN.CENTER)
    footer(s, "02:00"); notes(s, "[2 phút] Đi từ trái sang phải. Nhấn mạnh hybrid retrieval vì vector search đơn lẻ dễ bỏ lỡ mã chính sách/từ khóa chính xác. RBAC xảy ra trước generation.")

    # 6. Ingestion
    s = add_slide(prs); title(s, "Từ tài liệu thô đến tri thức truy xuất được", "Ingestion pipeline", 6)
    steps = [
        ("1", "Extract", "PDF / DOCX / PPTX\nOCR khi cần"),
        ("2", "Normalize", "Markdown sạch\nheading hierarchy"),
        ("3", "Atomic chunk", "~420 ký tự\n60 ký tự overlap"),
        ("4", "Enrich", "version • scope\nstatus • role"),
        ("5", "Dual index", "Embedding + BM25\nChroma + JSON"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.72 + i * 2.52
        rect(s, x, 2.0, 2.05, 2.7, WHITE, True, RGBColor(211, 225, 236))
        rect(s, x + 0.68, 1.68, 0.7, 0.7, CYAN if i < 4 else GREEN, True)
        text(s, num, x + 0.68, 1.82, 0.7, 0.3, 17, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, head, x + 0.18, 2.65, 1.7, 0.35, 17, NAVY, True, align=PP_ALIGN.CENTER)
        text(s, body, x + 0.18, 3.3, 1.7, 0.75, 13, MUTED, align=PP_ALIGN.CENTER)
    text(s, "Vì sao chunk nhỏ?", 0.85, 5.25, 2.5, 0.35, 17, BLUE, True)
    text(s, "Giảm nhiễu context, bám đúng điều khoản và tăng khả năng trả lời câu hỏi cụ thể — nhưng vẫn giữ heading để chunk không mất nghĩa.", 3.0, 5.2, 9.2, 0.75, 15, INK)
    footer(s, "02:00"); notes(s, "[2 phút] Nêu trade-off: chunk nhỏ không có nghĩa cắt tùy tiện. Hệ thống chia theo heading/đoạn và lặp heading context để giữ ngữ nghĩa.")

    # 7. Query flow
    s = add_slide(prs); title(s, "Luồng xử lý một câu hỏi", "Runtime flow", 7)
    flow = ["Auth + RBAC", "Screen intent", "Rewrite query", "Hybrid retrieve", "Rerank + filter", "Generate", "Verify + handoff"]
    for i, label in enumerate(flow):
        x = 0.55 + i * 1.82
        color = NAVY if i < 2 else BLUE if i < 5 else GREEN
        rect(s, x, 2.15, 1.5, 1.15, color, True)
        text(s, label, x + 0.08, 2.43, 1.34, 0.55, 13, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1:
            text(s, "→", x + 1.5, 2.48, 0.32, 0.35, 17, MUTED, True, align=PP_ALIGN.CENTER)
    rect(s, 0.85, 4.1, 3.45, 1.25, PALE, True, RGBColor(207, 226, 240))
    text(s, "ALLOW", 1.08, 4.35, 0.8, 0.3, 14, GREEN, True)
    text(s, "Câu HR đủ rõ → chạy RAG", 1.08, 4.75, 2.7, 0.28, 12, MUTED)
    rect(s, 4.9, 4.1, 3.45, 1.25, PALE, True, RGBColor(207, 226, 240))
    text(s, "CLARIFY", 5.13, 4.35, 1.0, 0.3, 14, GOLD, True)
    text(s, "Thiếu dữ kiện → hỏi ngược", 5.13, 4.75, 2.7, 0.28, 12, MUTED)
    rect(s, 8.95, 4.1, 3.45, 1.25, PALE, True, RGBColor(207, 226, 240))
    text(s, "BLOCK", 9.18, 4.35, 0.8, 0.3, 14, RED, True)
    text(s, "Ngoài phạm vi / dữ liệu cấm", 9.18, 4.75, 2.7, 0.28, 12, MUTED)
    footer(s, "02:00"); notes(s, "[2 phút] Đây là slide kỹ thuật quan trọng. Giải thích query rewrite cho follow-up, nhưng các câu như ok/cảm ơn được xử lý ở conversation layer, không đưa vào RAG.")

    # 8. Conversation intelligence
    s = add_slide(prs); title(s, "Không chỉ trả lời — chatbot biết tiếp tục hội thoại", "Conversation logic", 8)
    rect(s, 0.8, 1.75, 5.7, 4.7, WHITE, True, RGBColor(211, 225, 236))
    text(s, "Hỏi ngược đúng lúc", 1.15, 2.08, 4.7, 0.4, 20, NAVY, True)
    bullet_list(s, [
        "Thiếu ngày nghỉ → hỏi khoảng thời gian.",
        "Có ngày → hỏi lý do và bàn giao.",
        "Đủ dữ liệu → dựng sẵn form để rà soát.",
        "Câu nhạy cảm → hướng chuyển tiếp HR.",
    ], 1.15, 2.75, 4.7, 2.5, 16)
    rect(s, 6.85, 1.75, 5.7, 4.7, NAVY, True)
    text(s, "Một session, hai giao diện", 7.2, 2.08, 4.7, 0.4, 20, GOLD, True)
    bullet_list(s, [
        "Chat widget và màn đầy đủ dùng cùng logic.",
        "Đồng bộ lịch sử theo tài khoản.",
        "Giữ trạng thái luồng tạo đơn khi chuyển màn.",
        "Hiểu ok, cảm ơn, hiểu rồi — không reset hội thoại.",
    ], 7.2, 2.75, 4.7, 2.5, 16, WHITE)
    footer(s, "01:30"); notes(s, "[1.5 phút] Đây là điểm khác biệt về UX. Có thể kể lỗi thực tế đã sửa: người dùng nhắn ok mà bot tự giới thiệu lại; hiện đã có conversation acknowledgement layer.")

    # 9. Security
    s = add_slide(prs); title(s, "Guardrail và bảo vệ dữ liệu", "Trust by design", 9)
    labels = [
        ("RBAC", "Lọc tài liệu theo role và department", BLUE),
        ("Privacy", "Không tiết lộ lương / hồ sơ người khác", RED),
        ("Prompt security", "Chặn system prompt, API key, token", NAVY),
        ("Scope control", "Chặn crypto, thời tiết, code, Wi-Fi", GOLD),
        ("Human handoff", "Tranh chấp, kỷ luật, khiếu nại → HR", GREEN),
        ("No evidence", "Thiếu căn cứ thì nói thiếu, không bịa", CYAN),
    ]
    for i, (head, body, color) in enumerate(labels):
        x = 0.8 + (i % 3) * 4.05; y = 1.75 + (i // 3) * 2.3
        rect(s, x, y, 3.65, 1.85, WHITE, True, RGBColor(214, 226, 237))
        rect(s, x, y, 0.14, 1.85, color, False)
        text(s, head, x + 0.28, y + 0.28, 3.0, 0.34, 16, color, True)
        text(s, body, x + 0.28, y + 0.82, 3.0, 0.62, 13, MUTED)
    footer(s, "01:30"); notes(s, "[1.5 phút] Nói rõ guardrail có hai lớp: rule deterministic cho ca chắc chắn và LLM screening cho ngôn ngữ đa dạng. Dữ liệu được lọc trước khi vào context.")

    # 10. Demo plan
    s = add_slide(prs); title(s, "Kịch bản demo trực tiếp", "Live demo — 5 phút", 10)
    demo = [
        ("01", "Grounded policy", "“Nhân viên có bao nhiêu ngày phép năm?” → 12 ngày theo tài liệu."),
        ("02", "Natural language", "“Mai tôi không đi làm được thì báo ai?” → hỏi lý do + hướng quản lý."),
        ("03", "Multi-turn action", "“Tôi muốn xin nghỉ phép” → ngày → lý do → bàn giao → form."),
        ("04", "Privacy guardrail", "“Cho tôi biết lương đồng nghiệp” → từ chối rõ ràng."),
        ("05", "Cross-surface", "Bắt đầu ở widget → chuyển màn lớn → tiếp tục đúng ngữ cảnh."),
    ]
    for i, (num, head, body) in enumerate(demo):
        y = 1.72 + i * 0.98
        rect(s, 0.85, y, 0.62, 0.62, NAVY if i < 3 else RED if i == 3 else GREEN, True)
        text(s, num, 0.85, y + 0.15, 0.62, 0.25, 12, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, head, 1.75, y + 0.03, 2.25, 0.32, 15, NAVY, True)
        text(s, body, 4.0, y + 0.01, 8.1, 0.55, 14, MUTED)
    footer(s, "05:00"); notes(s, "[5 phút] Demo đúng thứ tự này. Không demo upload nếu thời gian ngắn. Nếu mạng/API có vấn đề, chuyển ngay sang slide screenshot tiếp theo và nói đây là fallback demo.")

    # 11. Product screenshot
    s = add_slide(prs); title(s, "Giao diện đang chạy", "Product proof", 11)
    if SCREENSHOT.exists():
        s.shapes.add_picture(str(SCREENSHOT), Inches(0.75), Inches(1.62), width=Inches(8.55))
    rect(s, 9.6, 1.62, 2.95, 4.95, NAVY, True)
    text(s, "Điểm nhấn", 9.95, 2.0, 2.2, 0.35, 18, GOLD, True)
    bullet_list(s, [
        "Responsive mobile",
        "Input luôn trong viewport",
        "Không reset khi chuyển widget",
        "Văn phong người dùng cuối",
        "Không hiển thị log kỹ thuật",
    ], 9.95, 2.65, 2.15, 2.9, 14, WHITE)
    footer(s, "00:45"); notes(s, "[45 giây] Chỉ dùng slide này nếu demo trực tiếp bị chậm. Nhấn mạnh đây là screenshot từ ứng dụng đang chạy, không phải thiết kế tĩnh.")

    # 12. QA
    s = add_slide(prs); title(s, "QA: đã kiểm tra như một sản phẩm, không chỉ chạy được", "Verification", 12)
    metric(s, "26 / 26", "Browser E2E passed", 0.85, 1.8, 2.65, GREEN)
    metric(s, "14 / 14", "Focused RAG regression", 3.8, 1.8, 2.65, GREEN)
    metric(s, "320 px", "Narrow mobile verified", 6.75, 1.8, 2.65, CYAN)
    metric(s, "0", "Horizontal overflow", 9.7, 1.8, 2.65, BLUE)
    bullet_list(s, [
        "Nghiệp vụ: nghỉ phép, nghỉ ốm, chấm công, OT, bảo hiểm, onboarding, lương.",
        "An toàn: privacy, prompt injection, câu ngoài phạm vi, secret/access request.",
        "Hội thoại: ok, cảm ơn, hiểu rồi, hủy luồng, follow-up và đồng bộ hai giao diện.",
        "Frontend production build đã thành công.",
    ], 1.0, 3.65, 11.2, 2.2, 16)
    footer(s, "01:30"); notes(s, "[1.5 phút] Giải thích test chạy qua trình duyệt thật: login, mở chat, gõ câu hỏi, Enter và đọc phản hồi render. Không chỉ gọi function nội bộ.")

    # 13. Limitations
    s = add_slide(prs); title(s, "Những gì đã tốt — và giới hạn cần nói thẳng", "Readiness", 13)
    rect(s, 0.8, 1.75, 5.75, 4.8, RGBColor(235, 249, 243), True, RGBColor(188, 228, 208))
    text(s, "Sẵn sàng pilot nội bộ", 1.15, 2.08, 4.8, 0.4, 20, GREEN, True)
    bullet_list(s, [
        "Luồng chính ổn định và có regression test.",
        "Guardrail deterministic cho ca rủi ro cao.",
        "Responsive desktop/mobile.",
        "Không hard-code API key.",
    ], 1.15, 2.8, 4.7, 2.5, 16)
    rect(s, 6.85, 1.75, 5.7, 4.8, RGBColor(255, 247, 232), True, RGBColor(239, 215, 164))
    text(s, "Trước khi scale enterprise", 7.2, 2.08, 4.8, 0.4, 20, GOLD, True)
    bullet_list(s, [
        "Bổ sung telemetry và dashboard chất lượng.",
        "Golden dataset lớn hơn theo phòng ban.",
        "SSO / enterprise identity và audit retention.",
        "Quy trình versioning, approval và rollback tài liệu.",
    ], 7.2, 2.8, 4.7, 2.5, 16)
    footer(s, "01:30"); notes(s, "[1.5 phút] Nói thật: đủ dùng pilot nội bộ, chưa nên gọi là hoàn tất enterprise ở quy mô lớn nếu chưa có observability, SSO và governance tài liệu đầy đủ.")

    # 14. Roadmap
    s = add_slide(prs); title(s, "Roadmap đề xuất", "Next 6–8 weeks", 14)
    phases = [
        ("NOW", "Stabilize", "Mở pilot\nThu feedback\nTheo dõi lỗi", GREEN),
        ("+2W", "Measure", "Golden set\nLatency/cost\nAnswer rating", BLUE),
        ("+4W", "Govern", "Version policy\nApproval flow\nAudit dashboard", GOLD),
        ("+6W", "Integrate", "SSO\nHRM API\nSharePoint sync", NAVY),
    ]
    for i, (tag, head, body, color) in enumerate(phases):
        x = 0.85 + i * 3.08
        rect(s, x, 2.0, 2.68, 3.65, WHITE, True, RGBColor(211, 225, 236))
        rect(s, x, 2.0, 2.68, 0.16, color, False)
        text(s, tag, x + 0.2, 2.42, 2.25, 0.3, 12, color, True)
        text(s, head, x + 0.2, 2.95, 2.25, 0.4, 19, NAVY, True)
        text(s, body, x + 0.2, 3.75, 2.25, 1.1, 15, MUTED, align=PP_ALIGN.CENTER)
    footer(s, "01:00"); notes(s, "[1 phút] Đề xuất xin leader chốt pilot và tiêu chí đo. Không hứa tích hợp lớn ngay; ưu tiên đo chất lượng trước.")

    # 15. Technique trade-offs
    s = add_slide(prs); title(s, "Technique đã chọn: vì sao chọn và đánh đổi gì?", "Engineering decisions", 15)
    rows = [
        ("Gemini 2.5 Flash", "Nhanh, tiếng Việt tốt, chi phí phù hợp demo", "Phụ thuộc mạng/provider; cần fallback"),
        ("Dense + BM25", "Bắt được cả ngữ nghĩa và từ khóa/mã chính sách", "Pipeline phức tạp hơn, phải tune score"),
        ("Atomic chunks", "Context ít nhiễu, bám sát điều khoản", "Có thể mất ngữ cảnh; bù bằng heading + overlap"),
        ("RRF + rerank", "Hợp nhất hai ranking, tăng precision top-k", "Tăng latency và chi phí tính toán"),
        ("LangGraph", "Luồng allow/clarify/block rõ, dễ kiểm soát", "Nhiều state/node hơn code tuyến tính"),
        ("Chroma local", "Triển khai pilot nhanh, không cần hạ tầng riêng", "Chưa tối ưu cho scale và HA lớn"),
    ]
    rect(s, 0.7, 1.5, 11.9, 0.65, NAVY, True)
    text(s, "Technique", 0.8, 1.65, 2.3, 0.3, 12, WHITE, True)
    text(s, "Ưu điểm / lý do chọn", 3.15, 1.65, 4.3, 0.3, 12, WHITE, True)
    text(s, "Nhược điểm / cách kiểm soát", 7.65, 1.65, 4.8, 0.3, 12, WHITE, True)
    for i, (tech, pros, cons) in enumerate(rows):
        y = 2.25 + i * 0.7
        if i % 2 == 0:
            rect(s, 0.7, y - 0.08, 11.9, 0.65, PALE, False)
        text(s, tech, 0.85, y, 2.15, 0.42, 13, BLUE, True)
        text(s, pros, 3.15, y, 4.25, 0.48, 12, INK)
        text(s, cons, 7.65, y, 4.65, 0.48, 12, MUTED)
    footer(s, "02:00"); notes(s, "[2 phút] Không chỉ liệt kê công nghệ. Với từng lựa chọn, nói rõ bài toán nó giải quyết và cái giá phải trả. Điểm quan trọng: không có technique nào tốt tuyệt đối; các biện pháp giảm nhược điểm đã được đưa vào thiết kế.")

    # 16. Code walkthrough
    s = add_slide(prs); title(s, "Đọc code theo luồng dữ liệu", "Code walkthrough", 16)
    code_map = [
        ("frontend/src/components/\nHrPolicyChatPanel.jsx", "State hội thoại, widget/full chat, leave intake, mobile UX"),
        ("backend/main.py", "API contract, auth, query rewrite, chat history, handoff"),
        ("src/hr_policy_graph.py", "Orchestrate screen → retrieve → RBAC → validate → generate"),
        ("src/task4_chunking_indexing.py", "Heading-aware chunking, embedding và Chroma index"),
        ("src/task9_retrieval_pipeline.py", "Dense + BM25 → RRF → rerank → fallback"),
        ("src/task10_generation.py", "Prompt, intent, deterministic guardrail và fallback answer"),
    ]
    for i, (path, responsibility) in enumerate(code_map):
        x = 0.8 + (i % 2) * 6.0; y = 1.65 + (i // 2) * 1.65
        rect(s, x, y, 5.55, 1.3, WHITE, True, RGBColor(211, 225, 236))
        text(s, path, x + 0.22, y + 0.2, 2.25, 0.78, 12, BLUE, True, "Consolas")
        text(s, responsibility, x + 2.55, y + 0.2, 2.72, 0.78, 12, MUTED)
    rect(s, 0.8, 6.55, 11.55, 0.42, NAVY, True)
    text(s, "Cách debug: UI symptom → API response → graph step → retrieved chunks → generation rule", 1.0, 6.64, 11.1, 0.22, 11, WHITE, True, align=PP_ALIGN.CENTER)
    footer(s, "02:00"); notes(s, "[2 phút] Mở IDE nếu leader muốn. Đi theo một request từ component qua API đến graph và generation. Nhấn mạnh ranh giới trách nhiệm của từng file; không giải thích code theo kiểu đọc từng dòng.")

    # 17. Agent orchestration
    s = add_slide(prs); title(s, "Kiểm soát và điều phối coding agent", "Human-in-the-loop engineering", 17)
    stages = [
        ("1", "Đặt acceptance criteria", "Ví dụ: sửa đúng intent, không đụng chunking/RAG."),
        ("2", "Tái hiện lỗi", "Browser/API test trước khi sửa; lưu output làm bằng chứng."),
        ("3", "Khoanh module", "Đọc call flow và xác định owner của lỗi."),
        ("4", "Patch nhỏ", "Chỉ sửa file liên quan, không refactor lan rộng."),
        ("5", "Regression gate", "Unit → API → browser E2E → production build."),
        ("6", "Review kết quả", "So expected/actual; lỗi khác module chỉ báo cáo, không tự sửa."),
    ]
    for i, (num, head, body) in enumerate(stages):
        x = 0.75 + (i % 3) * 4.15; y = 1.7 + (i // 3) * 2.15
        rect(s, x, y, 3.75, 1.72, WHITE, True, RGBColor(211, 225, 236))
        rect(s, x + 0.18, y + 0.18, 0.52, 0.52, BLUE if i < 4 else GREEN, True)
        text(s, num, x + 0.18, y + 0.29, 0.52, 0.2, 11, WHITE, True, align=PP_ALIGN.CENTER)
        text(s, head, x + 0.85, y + 0.2, 2.55, 0.35, 14, NAVY, True)
        text(s, body, x + 0.85, y + 0.73, 2.55, 0.65, 12, MUTED)
    rect(s, 0.85, 6.05, 11.45, 0.62, PALE, True, RGBColor(207, 226, 240))
    text(s, "Agent hỗ trợ tốc độ; người làm dự án vẫn chịu trách nhiệm về scope, quyết định kỹ thuật và bằng chứng kiểm thử.", 1.1, 6.22, 10.95, 0.28, 13, NAVY, True, align=PP_ALIGN.CENTER)
    footer(s, "02:00"); notes(s, "[2 phút] Đưa ví dụ thật: lỗi ký tự đ được tìm bằng test privacy; lỗi ok reset do client intent; lỗi mobile được đo bằng viewport 320px. Mỗi lỗi chỉ sửa module sở hữu rồi chạy lại E2E. Coding agent không được tự mở rộng phạm vi.")

    # 18. Close
    s = add_slide(prs)
    rect(s, 0, 0, 13.333, 7.5, NAVY, False)
    text(s, "KẾT LUẬN", 0.9, 0.85, 2.0, 0.3, 12, GOLD, True)
    text(s, "Từ chatbot trả lời câu hỏi\nthành HR Helpdesk có kiểm soát.", 0.9, 1.55, 8.5, 1.35, 31, WHITE, True, "Aptos Display")
    text(s, "Đề xuất: mở pilot nội bộ có giới hạn, đo chất lượng theo golden set và phản hồi người dùng trước khi scale.", 0.95, 3.35, 8.8, 0.85, 18, RGBColor(205, 221, 238))
    rect(s, 0.95, 4.65, 3.2, 0.78, CYAN, True)
    text(s, "Q&A", 0.95, 4.86, 3.2, 0.3, 18, WHITE, True, align=PP_ALIGN.CENTER)
    text(s, "Demo • Architecture • Risk • Roadmap", 0.95, 6.6, 6.0, 0.3, 10, GOLD, True)
    notes(s, "[3–5 phút Q&A] Câu chốt: hệ thống đã sẵn sàng pilot nội bộ; bước tiếp theo không phải thêm nhiều feature mà là đo chất lượng, governance dữ liệu và tích hợp định danh.")

    # Keep the core story to 13 visible slides; detailed technical slides remain as appendix.
    order = [0, 3, 10, 1, 4, 14, 7, 15, 16, 9, 11, 12, 17, 2, 5, 6, 8, 13]
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        prs.slides._sldIdLst.remove(slide_id)
    for index in order:
        prs.slides._sldIdLst.append(slide_ids[index])
    for slide in list(prs.slides)[13:]:
        slide._element.set("show", "0")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"slides={len(prs.slides)}")


if __name__ == "__main__":
    build()
